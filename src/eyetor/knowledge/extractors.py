"""Document extractors: bytes → plain text + lightweight section metadata.

Each extractor returns an ``ExtractedDoc`` with raw ``text`` and optionally a
list of ``sections`` (start/end offsets + heading path) that the chunker can
honour. Binary extractors (pdf/docx/xlsx/pptx) use lazy imports so the core
package keeps zero hard dependencies — if a library is missing, the extractor
logs a warning and returns ``None``.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable

logger = logging.getLogger(__name__)

_MISSING_DEP_WARNED: set[str] = set()


@dataclass
class ExtractedSection:
    """A logical section inside a document (start/end line, heading path)."""

    heading_path: str
    start_line: int
    end_line: int


@dataclass
class ExtractedDoc:
    """Normalized representation of a document for the chunker."""

    text: str
    title: str | None = None
    sections: list[ExtractedSection] = field(default_factory=list)


def _warn_missing(dep: str, ext: str) -> None:
    if dep in _MISSING_DEP_WARNED:
        return
    _MISSING_DEP_WARNED.add(dep)
    logger.warning(
        "extractor for %s unavailable: %s not installed (install eyetor[knowledge])",
        ext,
        dep,
    )


# ---------------------------------------------------------------------------
# Text / Markdown
# ---------------------------------------------------------------------------


_TEXT_EXTS = {
    ".md",
    ".mdx",
    ".markdown",
    ".rst",
    ".txt",
    ".log",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".go",
    ".rs",
    ".java",
    ".rb",
    ".php",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".cs",
    ".swift",
    ".kt",
    ".yaml",
    ".yml",
    ".toml",
    ".json",
    ".ini",
    ".cfg",
    ".env",
    ".sh",
    ".bash",
    ".zsh",
    ".sql",
    ".html",
    ".css",
    ".scss",
}


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        try:
            return path.read_text(encoding="latin-1")
        except Exception as exc:
            logger.debug("Failed to read %s: %s", path, exc)
            return ""
    except OSError as exc:
        logger.debug("Failed to open %s: %s", path, exc)
        return ""


def extract_text(path: Path) -> ExtractedDoc | None:
    text = _read_text(path)
    if not text.strip():
        return None
    title = _first_heading(text) or path.stem
    sections = _parse_markdown_sections(text) if path.suffix.lower() in {".md", ".mdx", ".markdown"} else []
    return ExtractedDoc(text=text, title=title, sections=sections)


def _first_heading(text: str) -> str | None:
    for line in text.splitlines()[:20]:
        stripped = line.strip()
        if stripped.startswith("#"):
            return stripped.lstrip("#").strip() or None
    return None


def _parse_markdown_sections(text: str) -> list[ExtractedSection]:
    """Walk markdown headings to build a hierarchical heading_path per section."""
    sections: list[ExtractedSection] = []
    stack: list[tuple[int, str]] = []  # (level, title)
    current_start = 0
    current_path: str | None = None
    lines = text.splitlines()
    for i, line in enumerate(lines):
        stripped = line.lstrip()
        if not stripped.startswith("#"):
            continue
        level = 0
        while level < len(stripped) and stripped[level] == "#":
            level += 1
        if level == 0 or level > 6:
            continue
        title = stripped[level:].strip()
        if not title:
            continue
        if current_path is not None and i > current_start:
            sections.append(
                ExtractedSection(
                    heading_path=current_path,
                    start_line=current_start,
                    end_line=i - 1,
                )
            )
        while stack and stack[-1][0] >= level:
            stack.pop()
        stack.append((level, title))
        current_path = " > ".join(t for _, t in stack)
        current_start = i
    if current_path is not None and current_start < len(lines):
        sections.append(
            ExtractedSection(
                heading_path=current_path,
                start_line=current_start,
                end_line=len(lines) - 1,
            )
        )
    return sections


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def extract_pdf(path: Path) -> ExtractedDoc | None:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        _warn_missing("pypdf", ".pdf")
        return None
    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        logger.warning("Failed to open PDF %s: %s", path, exc)
        return None
    parts: list[str] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            logger.debug("PDF page %d extract failed: %s", idx, exc)
            page_text = ""
        if page_text.strip():
            parts.append(f"\n\n[Page {idx}]\n\n{page_text}")
    text = "".join(parts).strip()
    if not text:
        return None
    title = None
    try:
        meta = reader.metadata or {}
        title = (meta.get("/Title") or None) if isinstance(meta, dict) else getattr(meta, "title", None)
    except Exception:
        pass
    return ExtractedDoc(text=text, title=title or path.stem)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def extract_docx(path: Path) -> ExtractedDoc | None:
    try:
        import docx  # type: ignore  # python-docx
    except ImportError:
        _warn_missing("python-docx", ".docx")
        return None
    try:
        document = docx.Document(str(path))
    except Exception as exc:
        logger.warning("Failed to open DOCX %s: %s", path, exc)
        return None

    lines: list[str] = []
    sections: list[ExtractedSection] = []
    stack: list[tuple[int, str]] = []
    current_path: str | None = None
    current_start = 0

    def flush_section(end_line: int) -> None:
        nonlocal current_path, current_start
        if current_path is not None and end_line >= current_start:
            sections.append(
                ExtractedSection(
                    heading_path=current_path,
                    start_line=current_start,
                    end_line=end_line,
                )
            )

    for para in document.paragraphs:
        text = (para.text or "").strip()
        style = (para.style.name if para.style else "") or ""
        if style.startswith("Heading"):
            try:
                level = int(style.split()[-1])
            except (ValueError, IndexError):
                level = 1
            flush_section(len(lines) - 1)
            while stack and stack[-1][0] >= level:
                stack.pop()
            stack.append((level, text or f"Heading {level}"))
            current_path = " > ".join(t for _, t in stack)
            current_start = len(lines)
            lines.append(("#" * level) + " " + (text or ""))
        elif text:
            lines.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            lines.append("\t".join(cells))

    flush_section(len(lines) - 1)
    body = "\n".join(lines).strip()
    if not body:
        return None
    return ExtractedDoc(text=body, title=path.stem, sections=sections)


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def extract_xlsx(path: Path) -> ExtractedDoc | None:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        _warn_missing("openpyxl", ".xlsx")
        return None
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        logger.warning("Failed to open XLSX %s: %s", path, exc)
        return None

    parts: list[str] = []
    sections: list[ExtractedSection] = []
    cursor = 0
    for sheet in wb.worksheets:
        header_line = f"# Sheet: {sheet.title}"
        start = cursor
        parts.append(header_line)
        cursor += 1
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if not any(cells):
                continue
            parts.append("\t".join(cells))
            cursor += 1
        if cursor > start:
            sections.append(
                ExtractedSection(
                    heading_path=sheet.title,
                    start_line=start,
                    end_line=cursor - 1,
                )
            )
    try:
        wb.close()
    except Exception:
        pass
    text = "\n".join(parts).strip()
    if not text:
        return None
    return ExtractedDoc(text=text, title=path.stem, sections=sections)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def extract_pptx(path: Path) -> ExtractedDoc | None:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        _warn_missing("python-pptx", ".pptx")
        return None
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        logger.warning("Failed to open PPTX %s: %s", path, exc)
        return None

    parts: list[str] = []
    sections: list[ExtractedSection] = []
    cursor = 0
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            pass
        heading_path = f"Slide {idx}" + (f": {title}" if title else "")
        start = cursor
        parts.append(f"# {heading_path}")
        cursor += 1
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    parts.append(text)
                    cursor += 1
        try:
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        except Exception:
            notes = ""
        if notes and notes.strip():
            parts.append(f"[Notes] {notes.strip()}")
            cursor += 1
        sections.append(
            ExtractedSection(
                heading_path=heading_path,
                start_line=start,
                end_line=cursor - 1,
            )
        )
    text = "\n".join(parts).strip()
    if not text:
        return None
    return ExtractedDoc(text=text, title=path.stem, sections=sections)


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------


Extractor = Callable[[Path], "ExtractedDoc | None"]


_EXTRACTORS: dict[str, Extractor] = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}


def get_extractor(suffix: str) -> Extractor | None:
    suffix = suffix.lower()
    if suffix in _EXTRACTORS:
        return _EXTRACTORS[suffix]
    if suffix in _TEXT_EXTS:
        return extract_text
    return None


def register_extractor(suffix: str, fn: Extractor) -> None:
    _EXTRACTORS[suffix.lower()] = fn


def supported_extensions() -> list[str]:
    return sorted(set(_TEXT_EXTS) | set(_EXTRACTORS.keys()))
